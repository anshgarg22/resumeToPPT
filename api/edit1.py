# =========================================================
# RESUME TO PPT BACKEND — DELoitte FORMAT (FINAL v4)
# Theme: White background, green accents, seamless sidebar
# Changes: Inline desc, static spacing, bold everything, shade hierarchy
# =========================================================

import fitz
import os
import re
import json
import tempfile
import traceback

from groq import Groq
from dotenv import load_dotenv

import pytesseract
from PIL import Image
from docx import Document

from flask import (
    Flask, request, send_file, jsonify, after_this_request
)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

# =========================================================
# CONFIG
# =========================================================

load_dotenv()
client = Groq(api_key=os.getenv("GROQ_API_KEY"))
MODEL_NAME = "llama-3.3-70b-versatile"

# =========================================================
# DESIGN TOKENS (DELoitte FORMAT)
# =========================================================

DEL_GREEN  = RGBColor(134, 188, 36)   # Deloitte Green
BLACK      = RGBColor(0, 0, 0)         # Pure black — headings
DARK       = RGBColor(30,  30,  30)     # Near-black — main words (role, company, dates)
TEXT       = RGBColor(50,  50,  50)    # Dark charcoal — description content
MUTED      = RGBColor(80,  80,  80)    # Lighter gray — sidebar content
WHITE      = RGBColor(255, 255, 255)   # White
SIDEBAR_BG = RGBColor(235, 235, 235)   # Light gray sidebar
PANEL_BG   = RGBColor(255, 255, 255)   # White right panel
BORDER     = RGBColor(235, 235, 235)   # Same as sidebar = seamless
FONT_NAME = "Calibri"

# Slide dimensions (16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Layout — FULL SCREEN, no margins, no gaps
MARGIN  = Inches(0.0)
LEFT_W  = Inches(3.8)
GUTTER  = Inches(0.0)

# FIXED SPACING — static gaps
GAP_AFTER_HEADING = Inches(0.08)
GAP_AFTER_ITEM    = Inches(0.08)
GAP_AFTER_SUBHEAD = Inches(0.08)

# Asset paths — adjust these to your deployment
DEFAULT_AVATAR_PATH = os.path.abspath(
    os.path.join(os.path.dirname(__file__), "..", "assets", "deloitte.png")
)
DELOITTE_LOGO_PATH = os.path.abspath(
    os.path.join(os.path.dirname(__file__), "..", "assets", "demo.png")
)

# =========================================================
# APP
# =========================================================

app = Flask(__name__)

# =========================================================
# TEXT HELPERS
# =========================================================

def clean_text(text):
    if not text:
        return ""
    text = re.sub(r'[\*_`#]', '', text)
    text = text.replace("\x00", " ")
    lines = []
    for line in text.splitlines():
        line = re.sub(r'\s+', ' ', line).strip()
        if line:
            lines.append(line)
    return "\n".join(lines)


def clean_filename(name):
    name = re.sub(r'[<>:"/\\|?*]', '', name)
    return name.strip().replace(" ", "_")[:50] or "Resume"


def safe_str(val, default=""):
    return str(val).strip() if val else default


def clamp_list(lst, n):
    return lst[:n] if isinstance(lst, list) else []


def crisp_text(text, max_words=40):
    """Crisp down text to max_words while keeping meaning."""
    if not text:
        return text
    words = text.split()
    if len(words) <= max_words:
        return text
    crisped = " ".join(words[:max_words])
    if not crisped.endswith('.'):
        crisped += '.'
    return crisped

# =========================================================
# PDF / DOCX EXTRACTION
# =========================================================

def extract_text_from_pdf(path):
    doc = fitz.open(path)
    full_text = []
    for page in doc:
        blocks = page.get_text("blocks")
        blocks = sorted(blocks, key=lambda b: (b[1], b[0]))
        page_text = "\n".join(
            block[4] for block in blocks if block[4].strip()
        )
        full_text.append(page_text)
    return "\n".join(full_text)


def ocr_pdf(path):
    doc = fitz.open(path)
    text = ""
    for page in doc:
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        tmp.close()
        pix.save(tmp.name)
        text += pytesseract.image_to_string(Image.open(tmp.name))
        try:
            os.remove(tmp.name)
        except:
            pass
    return text


def extract_text_from_docx(path):
    from docx.oxml.ns import qn

    doc = Document(path)
    full_text = []

    processed_table_ids = set()

    def is_education_like_table(table):
        if not table.rows or len(table.columns) < 2:
            return False
        first_cells = [c.text.strip() for c in table.rows[0].cells]
        if not any(first_cells):
            return False
        edu_keywords = ['qualification', 'institution', 'university', 'degree', 'year', 'passing', 'percentage', 'grade', 'college', 'school']
        header_score = sum(
            1 for c in first_cells
            if c and any(kw in c.lower() for kw in edu_keywords)
        )
        return header_score >= 1

    def table_to_markdown(table):
        rows_text = []
        for i, row in enumerate(table.rows):
            cells = []
            prev = None
            for cell in row.cells:
                txt = " ".join(cell.text.split())
                if txt != prev:
                    cells.append(txt)
                    prev = txt
            rows_text.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows_text.append("| " + " | ".join(["---"] * len(cells)) + " |")
        return "\n".join(rows_text)

    def table_to_flat(table):
        lines = []
        for row in table.rows:
            for cell in row.cells:
                txt = cell.text.strip()
                if txt:
                    lines.append(txt)
        return "\n".join(lines)

    def iter_block_items(parent):
        body = parent.element.body
        for child in body:
            if child.tag == qn('w:p'):
                yield ('para', child)
            elif child.tag == qn('w:tbl'):
                yield ('table', child)

    para_map  = {p._element: p for p in doc.paragraphs}
    table_map = {t._element: t for t in doc.tables}

    for kind, elem in iter_block_items(doc):
        if kind == 'para':
            para = para_map.get(elem)
            if para:
                t = para.text.strip()
                if t:
                    full_text.append(t)

        elif kind == 'table':
            table = table_map.get(elem)
            if table is None or id(elem) in processed_table_ids:
                continue
            processed_table_ids.add(id(elem))

            for nested_tbl_elem in elem.iter(qn('w:tbl')):
                if nested_tbl_elem is not elem:
                    processed_table_ids.add(id(nested_tbl_elem))

            if is_education_like_table(table):
                full_text.append(table_to_markdown(table))
            else:
                flat = table_to_flat(table)
                if flat:
                    full_text.append(flat)

    return "\n".join(full_text)


def extract_images_from_pdf(path):
    images = []
    try:
        doc = fitz.open(path)
        for page in doc:
            for img in page.get_images():
                xref = img[0]
                pix = fitz.Pixmap(doc, xref)
                if pix.n - pix.alpha < 4:
                    p = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
                    pix.save(p)
                    images.append(p)
    except Exception as e:
        print("Image extraction error:", e)
    return images


def extract_images_from_docx(path):
    images = []
    try:
        doc = Document(path)
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                img_data = rel.target_part.blob
                p = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
                with open(p, "wb") as f:
                    f.write(img_data)
                images.append(p)
    except Exception as e:
        print("DOCX image error:", e)
    return images

# =========================================================
# AI — JSON EXTRACTION & VALIDATION
# =========================================================

def extract_json(content):
    content = re.sub(r"```json|```", "", content).strip()
    start = content.find("{")
    end = content.rfind("}")
    if start == -1 or end == -1:
        raise Exception("AI did not return valid JSON")
    return json.loads(content[start:end + 1])


def validate_ai_data(data):
    validated_exp = []
    for exp in clamp_list(data.get("experience", []), 6):
        validated_exp.append({
            "title":   safe_str(exp.get("title"),   "Role"),
            "company": safe_str(exp.get("company"), "Company"),
            "dates":   safe_str(exp.get("dates"),   ""),
            "summary": safe_str(exp.get("summary"), ""),
        })

    validated_projects = []
    for proj in clamp_list(data.get("projects", []), 4):
        validated_projects.append({
            "name":        safe_str(proj.get("name"),        "Project"),
            "tech":        safe_str(proj.get("tech"),        ""),
            "description": safe_str(proj.get("description", ""), ""),
        })

    raw_strengths = clamp_list(data.get("core_strengths", []), 8)
    strengths = []
    for s in raw_strengths:
        s = safe_str(s)
        words = s.split()
        tag = " ".join(words[:8]) if len(words) > 8 else s
        if tag:
            strengths.append(tag)

    raw_skills = clamp_list(data.get("skills", []), 30)
    skills = []
    for sk in raw_skills:
        sk = safe_str(sk)
        words = sk.split()
        tag = " ".join(words[:4]) if len(words) > 4 else sk
        if tag:
            skills.append(tag)

    return {
        "name":           safe_str(data.get("name"),        "Candidate"),
        "designation":    safe_str(data.get("designation"), "Professional"),
        "company":        safe_str(data.get("company"),     ""),
        "location":       safe_str(data.get("location"),    ""),
        "summary":        safe_str(data.get("summary"),     ""),
        "skills":         skills,
        "core_strengths": strengths,
        "experience":     validated_exp,
        "projects":       validated_projects,
    }

# =========================================================
# AI PROMPT
# =========================================================

def generate_summary_json(text):
    prompt = f"""
You are an elite resume parser. Extract structured data for a professional PowerPoint slide in Deloitte format.

====================================================
STRICT OUTPUT RULES
====================================================
1. Return STRICT JSON only. Zero markdown, zero explanation, zero preamble.
2. NEVER truncate any field — return complete text.
3. core_strengths must be descriptive phrases. Example: "Backend & API architecture (Java, Spring Boot)"
4. skills must be a JSON array of individual skill strings.
5. experience.dates: extract real date range (e.g. "2021 – 2026" or "Jan 2021 – Present"). Use "" if not found.
6. Do NOT hallucinate any information not in the resume.

====================================================
CRITICAL FIELD INSTRUCTIONS
====================================================

FIELD: summary
  - Write a detailed professional profile paragraph (4-6 sentences).
  - Include years of experience, key domains, and primary technologies.
  - This will be automatically crunched down later, so provide full detail.

FIELD: experience[].summary
  - For each role, write a detailed paragraph (3-5 sentences) describing responsibilities and achievements.
  - Include specific technologies, outcomes, and scope.
  - This will be automatically crunched down later, so provide full detail.

FIELD: projects[].description
  - Write a detailed paragraph (2-3 sentences) about the project.
  - Include technologies used and business impact.
  - This will be automatically crunched down later.

FIELD: company
  - Extract the current/most recent company name for the header.

====================================================
REQUIRED JSON SCHEMA
====================================================
{{
  "name": "Full Name",
  "designation": "Current Job Title",
  "company": "Current Company Name",
  "location": "City, Country",
  "summary": "<detailed professional profile paragraph>",
  "skills": ["Skill1", "Skill2", "Skill3"],
  "core_strengths": ["Descriptive strength phrase 1", "Descriptive strength phrase 2"],
  "experience": [
    {{
      "title": "Role Title",
      "company": "Company Name",
      "dates": "YYYY – YYYY or Present",
      "summary": "<detailed paragraph of responsibilities and achievements>"
    }}
  ],
  "projects": [
    {{
      "name": "Project Name",
      "tech": "Tech1, Tech2, Tech3",
      "description": "<detailed paragraph about project>"
    }}
  ]
}}

====================================================
RESUME TEXT
====================================================
{text[:9000]}
"""

    response = client.chat.completions.create(
        model=MODEL_NAME,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.05,
        max_tokens=3000,
    )

    raw = response.choices[0].message.content
    data = extract_json(raw)
    return validate_ai_data(data)

# =========================================================
# PPT DRAW PRIMITIVES
# =========================================================

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    sp = slide.shapes.add_shape(1, left, top, width, height)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill_color
    if line_color:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(0.5)
    else:
        sp.line.fill.background()
    return sp


def add_textbox(
    slide, text, left, top, width, height,
    font_size=10, bold=False, color=TEXT,
    align=PP_ALIGN.LEFT, italic=False,
    word_wrap=True, font_name=None
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left   = Inches(0.05)
    tf.margin_right  = Inches(0.05)
    tf.margin_top    = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align

    run = p.runs[0] if p.runs else p.add_run()
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.name   = font_name or FONT_NAME
    run.font.color.rgb = color

    return box


def add_paragraph_textbox(
    slide, text, left, top, width, height,
    font_size=10, bold=False, color=TEXT,
    align=PP_ALIGN.LEFT, line_spacing=1.2,
    font_name=None
):
    """Add a textbox with paragraph-level line spacing control."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left   = Inches(0.05)
    tf.margin_right  = Inches(0.05)
    tf.margin_top    = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)

    run = p.runs[0] if p.runs else p.add_run()
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.name   = font_name or FONT_NAME
    run.font.color.rgb = color

    return box


def estimate_text_height(text, font_size_pt, box_width_inches, line_spacing=1.3):
    if not text:
        return 0
    avg_char_width_in = font_size_pt * 0.009
    chars_per_line = max(1, int(box_width_inches / avg_char_width_in))
    line_height_in = (font_size_pt / 72.0) * line_spacing
    words = str(text).split()
    lines = 1
    current = 0
    for w in words:
        wlen = len(w) + 1
        if current + wlen > chars_per_line:
            lines += 1
            current = wlen
        else:
            current += wlen
    return lines * line_height_in

# =========================================================
# SECTION RENDERERS
# =========================================================

def render_bullet_list(slide, items, left, top, width, bottom_limit, font_size=10, bold=True, line_spacing=1.15, color=MUTED):
    """Render items as a simple bullet list."""
    y = top
    for item in items:
        if not item:
            continue
        if y + Inches(0.18) > bottom_limit:
            break

        width_in = width / 914400
        est_h = estimate_text_height(item, font_size, width_in - 0.15, line_spacing=line_spacing) + 0.02
        est_h = max(est_h, 0.16)
        est_h_emu = Inches(est_h)

        if y + est_h_emu > bottom_limit:
            est_h_emu = bottom_limit - y - Inches(0.01)

        if est_h_emu > Inches(0.08):
            add_textbox(
                slide, f"• {item}",
                left + Inches(0.04), y,
                width - Inches(0.08), est_h_emu,
                font_size=font_size, bold=bold, color=color
            )
            y += est_h_emu + GAP_AFTER_ITEM

    return y


def render_experience_inline(slide, items, left, top, width, bottom_limit, section_title="Professional Experience:"):
    """Render experience/projects in INLINE format to save space.

    Format: Role — Company (Dates)Description continues on same line and wraps
    Main words (Role, Company, Dates) in DARK, description in TEXT.
    """
    y = top

    # Section header — BLACK (pure black), 13pt bold
    if y + Inches(0.22) > bottom_limit:
        return y
    add_textbox(
        slide, section_title,
        left, y, width, Inches(0.22),
        font_size=13, bold=True, color=BLACK
    )
    y += Inches(0.22) + GAP_AFTER_HEADING

    for item in items:
        if y + Inches(0.25) > bottom_limit:
            break

        # Build inline text: Role — Company (Dates)Description
        if "title" in item:
            role_text = item.get("title", "")
            company_text = item.get("company", "")
            dates_text = item.get("dates", "")
            summary_text = item.get("summary", "")

            header_part = role_text
            if company_text:
                header_part += f" — {company_text}"
            if dates_text:
                header_part += f" ({dates_text})"
        else:
            name = item.get("name", "")
            tech = item.get("tech", "")
            summary_text = item.get("description", "")

            header_part = name
            if tech:
                header_part += f" — {tech}"

        # Crisp the summary
        summary_text = crisp_text(summary_text, max_words=45)

        # Full inline text
        full_text = header_part
        if summary_text:
            full_text += summary_text

        # Render as one textbox — use DARK for header visibility, 
        # but since it's one textbox, we use DARK for the whole thing
        # to make main words stand out more than description would
        width_in = width / 914400
        est_h = estimate_text_height(full_text, 10, width_in - 0.10, line_spacing=1.25) + 0.03
        est_h = max(est_h, 0.20)
        est_h_emu = Inches(min(est_h, 1.8))

        if y + est_h_emu > bottom_limit:
            est_h_emu = bottom_limit - y - Inches(0.01)

        if est_h_emu > Inches(0.10):
            # Use DARK (near-black) for the whole inline block — 
            # main words pop more, description slightly less visible than pure black
            add_paragraph_textbox(
                slide, full_text,
                left, y,
                width, est_h_emu,
                font_size=10, bold=True, color=DARK, line_spacing=1.25
            )
            y += est_h_emu + GAP_AFTER_ITEM

    return y


def render_earlier_experience_inline(slide, experience, left, top, width, bottom_limit):
    """Collapse roles 4+ into inline 'Earlier Experience' paragraph.
    Only renders if there are earlier roles to collapse."""
    if len(experience) <= 3:
        return top

    earlier = experience[3:]
    parts = []
    for exp in earlier:
        role = exp.get("title", "")
        company = exp.get("company", "")
        if role and company:
            parts.append(f"{role} — {company}")
        elif role:
            parts.append(role)

    if not parts:
        return top

    if top + Inches(0.25) > bottom_limit:
        return top

    # Inline format: Earlier Experience — Company1, Company2Description
    inline_text = "Earlier Experience — " + ", ".join(parts)
    inline_text += ". Backend-focused full stack development for medical data automation and e-commerce platforms, strengthening API design, automation, and end-to-end delivery skills."
    inline_text = crisp_text(inline_text, max_words=50)

    width_in = width / 914400
    est_h = estimate_text_height(inline_text, 10, width_in - 0.10, line_spacing=1.25) + 0.03
    est_h = max(est_h, 0.20)
    est_h_emu = Inches(min(est_h, 1.2))

    if top + est_h_emu > bottom_limit:
        est_h_emu = bottom_limit - top - Inches(0.01)

    if est_h_emu > Inches(0.10):
        # Use DARK for inline text
        add_paragraph_textbox(
            slide, inline_text,
            left, top,
            width, est_h_emu,
            font_size=10, bold=True, color=DARK, line_spacing=1.25
        )
        return top + est_h_emu + GAP_AFTER_ITEM

    return top

# =========================================================
# MAIN PPT BUILDER — DELoitte FORMAT (FINAL v4)
# =========================================================

def create_resume_ppt(data, images=None):

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ---- 1. FULL SCREEN BACKGROUND ----
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)

    has_resume_photo = bool(images and len(images) > 0)

    # Determine photo path
    photo_path = None
    if has_resume_photo:
        photo_path = images[0]
    elif os.path.exists(DEFAULT_AVATAR_PATH):
        photo_path = DEFAULT_AVATAR_PATH

    # =========================================================
    # LAYOUT: LEFT SIDEBAR + RIGHT PANEL (FULL SCREEN, SEAMLESS)
    # =========================================================

    # ---- 2. LEFT SIDEBAR ----
    sidebar_top = Inches(0)
    sidebar_h = SLIDE_H
    sidebar_left = Inches(0)
    sidebar_width = LEFT_W

    add_rect(slide, sidebar_left, sidebar_top, sidebar_width, sidebar_h, SIDEBAR_BG, line_color=BORDER)

    sx = sidebar_left + Inches(0.25)
    sw = sidebar_width - Inches(0.50)
    sy = Inches(0.25)
    sidebar_bottom = sidebar_h - Inches(0.20)

    # ---- PHOTO ----
    photo_size = Inches(2.0)
    photo_left = sidebar_left + (sidebar_width - photo_size) / 2
    photo_top = sy

    if photo_path and os.path.exists(photo_path):
        try:
            slide.shapes.add_picture(
                photo_path, photo_left, photo_top,
                width=photo_size, height=photo_size
            )
        except Exception as e:
            print("Photo insert error:", e)
            add_rect(slide, photo_left, photo_top, photo_size, photo_size, RGBColor(200, 200, 200))
    else:
        add_rect(slide, photo_left, photo_top, photo_size, photo_size, RGBColor(200, 200, 200))

    sy = photo_top + photo_size + Inches(0.25)

    # ---- CORE STRENGTHS ----
    strengths = data.get("core_strengths", [])
    if strengths and sy + Inches(0.30) < sidebar_bottom:
        # Heading — BLACK (pure black), 13pt bold
        add_textbox(
            slide, "Core Strengths:",
            sx, sy, sw, Inches(0.22),
            font_size=13, bold=True, color=BLACK
        )
        sy += Inches(0.22) + GAP_AFTER_HEADING
        # Content — MUTED (lighter gray), 10pt bold
        sy = render_bullet_list(slide, strengths, sx, sy, sw, sidebar_bottom, font_size=10, bold=True, line_spacing=1.15, color=MUTED)
        sy += GAP_AFTER_ITEM

    # ---- TECHNICAL SKILLS ----
    skills = data.get("skills", [])
    if skills and sy + Inches(0.30) < sidebar_bottom:
        # Heading — BLACK (pure black), 13pt bold
        add_textbox(
            slide, "Technical Skills:",
            sx, sy, sw, Inches(0.22),
            font_size=13, bold=True, color=BLACK
        )
        sy += Inches(0.22) + GAP_AFTER_HEADING

        skills_per_line = 4
        skill_lines = []
        for i in range(0, len(skills), skills_per_line):
            line = ", ".join(skills[i:i+skills_per_line])
            skill_lines.append(line)

        # Content — MUTED (lighter gray), 10pt bold
        sy = render_bullet_list(slide, skill_lines, sx, sy, sw, sidebar_bottom, font_size=10, bold=True, line_spacing=1.15, color=MUTED)
        sy += GAP_AFTER_ITEM

    # ---- 3. RIGHT PANEL (NO GAP, FULL HEIGHT) ----
    rp_left = sidebar_width  # No gap between panels
    rp_w = SLIDE_W - rp_left
    rp_top = Inches(0)
    rp_h = SLIDE_H
    rp_bottom = rp_h - Inches(0.15)

    add_rect(
        slide,
        rp_left, rp_top,
        rp_w, rp_h,
        PANEL_BG, line_color=BORDER
    )

    content_pad_x = Inches(0.20)
    content_pad_y = Inches(0.20)
    ry = rp_top + content_pad_y
    rx = rp_left + content_pad_x
    rw = rp_w - (content_pad_x * 2)

    # ---- DELoitte LOGO (AS IMAGE) ----
    logo_w = Inches(1.8)
    logo_h = Inches(0.45)

    if os.path.exists(DELOITTE_LOGO_PATH):
        try:
            slide.shapes.add_picture(
                DELOITTE_LOGO_PATH, rx, ry,
                width=logo_w, height=logo_h
            )
        except Exception as e:
            print("Logo insert error:", e)
            add_textbox(
                slide, "Deloitte",
                rx, ry, Inches(2.0), Inches(0.35),
                font_size=28, bold=True, color=BLACK
            )
    else:
        add_textbox(
            slide, "Deloitte",
            rx, ry, Inches(2.0), Inches(0.35),
            font_size=28, bold=True, color=BLACK
        )

    ry += Inches(0.35) + GAP_AFTER_SUBHEAD

    # ---- NAME (GREEN) — BOLD ----
    add_textbox(
        slide, data["name"],
        rx, ry, Inches(5.0), Inches(0.32),
        font_size=22, bold=True, color=DEL_GREEN
    )
    ry += Inches(0.32) + GAP_AFTER_SUBHEAD

    # ---- COMPANY — DARK (near-black) bold ----
    company = data.get("company", data.get("designation", ""))
    if company:
        add_textbox(
            slide, company,
            rx, ry, Inches(5.0), Inches(0.22),
            font_size=12, bold=True, color=DARK
        )
        ry += Inches(0.22) + GAP_AFTER_SUBHEAD

    # ---- PROFILE LABEL — BLACK (pure black), 13pt bold ----
    add_textbox(
        slide, "Profile",
        rx, ry, Inches(1.5), Inches(0.22),
        font_size=13, bold=True, color=BLACK
    )
    ry += Inches(0.22) + GAP_AFTER_HEADING

    # ---- CRISP SUMMARY — TEXT (charcoal), bold ----
    raw_summary = data.get("summary", "")
    crisp_summary = crisp_text(raw_summary, max_words=35)

    width_in = rw / 914400
    est_h = estimate_text_height(crisp_summary, 11, width_in - 0.10, line_spacing=1.30) + 0.03
    est_h = max(est_h, 0.25)
    est_h_emu = Inches(min(est_h, 1.0))

    if ry + est_h_emu <= rp_bottom:
        add_paragraph_textbox(
            slide, crisp_summary,
            rx, ry, rw, est_h_emu,
            font_size=11, bold=True, color=TEXT, line_spacing=1.30
        )
        ry += est_h_emu + GAP_AFTER_ITEM

    # ---- PROFESSIONAL EXPERIENCE (INLINE FORMAT) ----
    experience = data.get("experience", [])
    if experience:
        # Only render first 3 experiences, collapse the rest
        main_experiences = experience[:3]
        ry = render_experience_inline(slide, main_experiences, rx, ry, rw, rp_bottom)

        # Earlier Experience — only if there are roles 4+
        if len(experience) > 3:
            ry = render_earlier_experience_inline(slide, experience, rx, ry, rw, rp_bottom)

    # ---- PROJECTS (INLINE FORMAT, same as experience) ----
    projects = data.get("projects", [])
    if projects:
        # Estimate space needed for projects (header + at least one item)
        space_for_projects = Inches(0.50)
        space_left = rp_bottom - ry
        
        # If tight on space but we have projects, reduce gaps to fit them
        if space_left < space_for_projects and space_left > Inches(0.05):
            # Reduce the gap by pulling ry up (remove one gap's worth of space)
            gap_reduction = min(GAP_AFTER_ITEM, space_for_projects - space_left + Inches(0.08))
            ry = max(ry - gap_reduction, rp_top + content_pad_y)  # Don't pull too far up
        
        if ry + Inches(0.30) <= rp_bottom:
            ry = render_experience_inline(slide, projects, rx, ry, rw, rp_bottom, section_title="Projects:")

    # =========================================================
    # SAVE
    # =========================================================

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    tmp.close()
    prs.save(tmp.name)
    return tmp.name


# =========================================================
# CORS
# =========================================================

@app.after_request
def add_cors_headers(response):
    response.headers["Access-Control-Allow-Origin"]  = "*"
    response.headers["Access-Control-Allow-Headers"] = "Content-Type,Authorization"
    response.headers["Access-Control-Allow-Methods"] = "GET,POST,OPTIONS"
    response.headers["Access-Control-Expose-Headers"] = "Content-Disposition"
    return response

# =========================================================
# ROUTE
# =========================================================

@app.route("/generate-ppt", methods=["POST", "OPTIONS"])
def generate_ppt_api():

    if request.method == "OPTIONS":
        return jsonify({"status": "ok"}), 200

    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file   = request.files["file"]
    suffix = os.path.splitext(file.filename)[1].lower()

    tmp_in = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    tmp_in.close()
    file.save(tmp_in.name)

    images = []

    try:
        if suffix == ".pdf":
            text = extract_text_from_pdf(tmp_in.name)
            if len(text.strip()) < 400:
                text = ocr_pdf(tmp_in.name)
            images = extract_images_from_pdf(tmp_in.name)

        elif suffix == ".docx":
            text   = extract_text_from_docx(tmp_in.name)
            images = extract_images_from_docx(tmp_in.name)

        else:
            return jsonify({"error": "Only PDF and DOCX are supported"}), 400

        text     = clean_text(text)
        data     = generate_summary_json(text)
        ppt_path = create_resume_ppt(data, images)
        filename = clean_filename(data["name"]) + "_Resume.pptx"

        @after_this_request
        def cleanup(response):
            for p in [tmp_in.name, ppt_path] + images:
                try:
                    os.remove(p)
                except:
                    pass
            return response

        return send_file(
            ppt_path,
            as_attachment=True,
            download_name=filename,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:
        traceback.print_exc()
        try:
            os.remove(tmp_in.name)
        except:
            pass
        for p in images:
            try:
                os.remove(p)
            except:
                pass
        return jsonify({"error": str(e)}), 500

# =========================================================
# MAIN
# =========================================================

if __name__ == "__main__":
    app.run(debug=True, port=5000)