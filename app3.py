import fitz
import os
import re
from groq import Groq
from dotenv import load_dotenv
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import platform

# -------- CONFIG --------
load_dotenv()
client = Groq(api_key=os.getenv("GROQ_API_KEY"))
MODEL_NAME = "llama-3.1-8b-instant"

# -------- THEME COLORS --------
PRIMARY = RGBColor(0, 102, 204)
TEXT_DARK = RGBColor(40, 40, 40)

# -------- PATH --------
def get_downloads_path():
    if platform.system() == "Windows":
        return os.path.join(os.environ["USERPROFILE"], "Downloads")
    return os.path.join(os.path.expanduser("~"), "Downloads")

# -------- CLEAN --------
def clean_text(text):
    text = re.sub(r'[\*\_`#]', '', text)
    return " ".join(text.replace("\n", " ").split())

def clean_filename(name):
    name = re.sub(r'[<>:"/\\|?*]', '', name)
    return name.strip().replace(" ", "_")[:50] or "Resume"

# -------- FILE READ --------
def extract_text_from_pdf(path):
    doc = fitz.open(path)
    return "".join([p.get_text() for p in doc])

def ocr_pdf(path):
    doc = fitz.open(path)
    text = ""
    for i, page in enumerate(doc):
        pix = page.get_pixmap()
        img_path = f"temp_{i}.png"
        pix.save(img_path)
        text += pytesseract.image_to_string(Image.open(img_path))
        os.remove(img_path)
    return text

def extract_text_from_docx(path):
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

# -------- GROQ --------
def generate_structured_output(text):

    prompt = f"""
You are a resume to PPT CONVERTER.

Generate structured output in this format :

🔹 Profile Overview
Name: ...
Experience: ...
Certifications: ...
Core Focus: ...
Domains Worked In: ...

🔹 Key Strengths
- ...
- ...

🔹 Technical Skills
Languages & Frameworks: ...
Databases: ...
Tools: ...
Cloud: ...

🔹 Professional Experience
🟢 Role – Company (Duration)
- ...                                                                                                                   
- ...
(change slide for each role)

🔹 Projects
- ...
- ...

🔹 Overall Impression
- ...
- ...

Rules:
- No markdown
- Keep it clean
- use bullets for points(start with "-")

Resume:
{text}
"""

    res = client.chat.completions.create(
        model=MODEL_NAME,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3
    )

    return res.choices[0].message.content

# -------- PARSE --------
def parse_sections(text):
    sections = {}
    current = None

    for line in text.split("\n"):
        line = line.strip()
        if not line:
            continue

        if line.startswith("🔹"):
            current = line.replace("🔹", "").strip()
            sections[current] = []
        else:
            if current:
                sections[current].append(line)

    return sections

# -------- PPT CREATION --------
def create_ppt(data):

    prs = Presentation()

    # -------- TITLE SLIDE --------
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    name = data.get("Profile Overview", ["Resume"])[0].replace("Name:", "").strip()

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Professional Resume Overview"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER

    # -------- CONTENT --------
    for section, lines in data.items():

        chunk = []
        for line in lines:
            chunk.append(line)

            if len(chunk) == 5:
                add_slide(prs, section, chunk)
                chunk = []

        if chunk:
            add_slide(prs, section, chunk)

    filename = f"{clean_filename(name)}_Resume.pptx"
    path = os.path.join(get_downloads_path(), filename)
    prs.save(path)

    return path

def add_slide(prs, title, lines):

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Left bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), prs.slide_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK

    # Underline
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.2), Inches(2), Inches(0.1)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = PRIMARY
    underline.line.fill.background()

    # Content
    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(8), Inches(5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line.replace("🟢", "").strip()

        if line.startswith("-"):
            p.level = 1
            p.font.size = Pt(18)
        else:
            p.level = 0
            p.font.size = Pt(20)

        p.font.color.rgb = TEXT_DARK

# -------- MAIN --------
# if __name__ == "__main__":

#     file_path = "sample2.pdf"  # change if needed

#     print("Reading file...")

#     if file_path.endswith(".pdf"):
#         text = extract_text_from_pdf(file_path)

#         if len(text.strip()) < 500:
#             print("Using OCR...")
#             text = ocr_pdf(file_path)

#     elif file_path.endswith(".docx"):
#         text = extract_text_from_docx(file_path)

#     else:
#         raise ValueError("Unsupported format")

#     text = clean_text(text)

#     print("Generating structured output...")
#     structured_output = generate_structured_output(text)

#     print("\nStructured Output:\n")
#     print(structured_output)

#     parsed = parse_sections(structured_output)

#     print("\nCreating PPT...")
#     path = create_ppt(parsed)

#     print(f"\nPPT saved at: {path}")


    # CONNECTING TO FRONTEND 
from flask import Flask, request, send_file
import tempfile

app = Flask(__name__)

@app.route("/generate-ppt", methods=["POST"])
def generate_ppt_api():
    file = request.files["file"]

    temp_input = tempfile.NamedTemporaryFile(delete=False)
    file.save(temp_input.name)

    # Extract text
    if file.filename.endswith(".pdf"):
        text = extract_text_from_pdf(temp_input.name)
        if len(text.strip()) < 500:
            text = ocr_pdf(temp_input.name)
    else:
        text = extract_text_from_docx(temp_input.name)

    text = clean_text(text)

    structured = generate_structured_output(text)
    parsed = parse_sections(structured)

    ppt_path = create_ppt(parsed)

    return send_file(
    ppt_path,
    as_attachment=True,
    download_name="resume_presentation.pptx",
    mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

if __name__ == "__main__":
    app.run(debug=True)